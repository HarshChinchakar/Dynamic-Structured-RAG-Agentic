# main.py — Streamlit front-end for Router -> (Mongo | Policy) handlers
import streamlit as st
from datetime import datetime
import traceback
import os

from src.Router_gpt import classify_query, RouteType
from src.Mutlimedia import (
    load_in_memory_vectorstore,
    get_cached_vectorstore,
    clear_cache,
    cache_status,
    policy_handler
)
from langchain_openai import OpenAIEmbeddings
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document
from langchain.text_splitter import CharacterTextSplitter

# ===========================================
# PAGE SETUP
# ===========================================
st.set_page_config(page_title="Mongo_RAG Streamlit", page_icon="📚", layout="wide")
st.title("📚 Mongo_RAG - Streamlit Frontend")

st.markdown("""
### 🧠 Intelligent Query Router + RAG Pipeline
The pipeline:
1. Upload policy documents (PDF, DOCX, PPTX)
2. Automatically embed & store them in-memory (Chroma)
3. Enter a natural-language query — the router decides whether it's:
   - **Policy** (RAG retrieval)
   - **Document** (Mongo query)
   - **Both** (Hybrid)
""")

# ===========================================
# FILE UPLOAD SECTION
# ===========================================
st.subheader("📂 Step 1 — Upload Policy Documents")

uploaded_files = st.file_uploader(
    "Upload one or more files (PDF, DOCX, PPTX)",
    type=["pdf", "docx", "pptx"],
    accept_multiple_files=True
)

if uploaded_files:
    with st.spinner("Processing and embedding uploaded files..."):
        try:
            all_texts = []
            splitter = CharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
            for uploaded_file in uploaded_files:
                file_ext = os.path.splitext(uploaded_file.name)[1].lower()
                text = ""

                # Extract text based on file type
                if file_ext == ".pdf":
                    reader = PdfReader(uploaded_file)
                    for page in reader.pages:
                        t = page.extract_text()
                        if t:
                            text += t + "\n"
                elif file_ext == ".pptx":
                    prs = Presentation(uploaded_file)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text += shape.text + "\n"
                elif file_ext == ".docx":
                    doc = Document(uploaded_file)
                    for para in doc.paragraphs:
                        text += para.text + "\n"

                chunks = splitter.split_text(text)
                all_texts.extend(chunks)

            if not all_texts:
                st.error("No text extracted from uploaded files.")
            else:
                if get_cached_vectorstore():
                    st.info("Existing in-memory Chroma DB found — clearing to avoid conflicts.")
                    clear_cache()

                embeddings = OpenAIEmbeddings()
                load_in_memory_vectorstore(embeddings, all_texts, collection_name="ram_store")
                st.success(f"✅ Embedded and cached {len(all_texts)} text chunks in memory.")
        except Exception as e:
            st.error(f"❌ Embedding process failed: {e}")
            st.code(traceback.format_exc())

st.caption(f"🧮 Cache status: {cache_status()}")

st.divider()

# ===========================================
# QUERY SECTION
# ===========================================
st.subheader("💬 Step 2 — Enter Your Query")
query = st.text_input("Enter your question", placeholder="e.g. What is the leave encashment policy?")

col1, col2 = st.columns([3, 1])
with col2:
    if st.button("Run Query"):
        if not query or not query.strip():
            st.warning("⚠️ Please enter a question.")
        else:
            st.session_state['run_time'] = datetime.utcnow().isoformat()
            st.session_state['query'] = query.strip()

if 'run_time' in st.session_state:
    st.caption(f"Last run UTC: {st.session_state['run_time']}")

# ===========================================
# ROUTING + EXECUTION PIPELINE
# ===========================================
if 'query' in st.session_state and st.session_state['query']:
    q = st.session_state['query']

    st.divider()
    st.subheader("🔍 Step 3 — Routing and Execution")

    with st.spinner("Classifying query via Router..."):
        try:
            router_response = classify_query(q)
        except Exception as e:
            st.error(f"❌ Router crashed: {e}")
            st.code(traceback.format_exc())
            st.stop()

    if isinstance(router_response, dict) and "error" in router_response:
        st.error(f"❌ Router error at stage **{router_response.get('stage')}**")
        st.code(router_response["error"])
        st.stop()
    else:
        try:
            route, confidence, reason, doc_q, pol_q = router_response
        except Exception as e:
            st.error(f"❌ Unexpected Router output format: {e}")
            st.write(router_response)
            st.stop()

    with st.expander("Router Output (Raw)", expanded=False):
        st.json({
            "route": getattr(route, "value", str(route)),
            "confidence": confidence,
            "reason": reason,
            "doc_query": doc_q,
            "policy_query": pol_q,
        })

    st.success(f"✅ Router classified as **{getattr(route, 'value', str(route)).upper()}** (confidence {confidence})")

    result_text = ""
    logs = []

    # -------------------- DOCUMENT HANDLER --------------------
    if getattr(route, "value", str(route)).lower() == "document":
        st.info("🗂️ Running Mongo (Document) handler...")
        logs.append("Router → Document Handler (Mongo)")
        try:
            from src.Mongo import query_mongo
            with st.spinner("Executing Mongo pipeline..."):
                result_text = query_mongo(q)
            if isinstance(result_text, str) and result_text.startswith("ERROR"):
                st.error(result_text)
            else:
                st.success("✅ Document handler executed successfully.")
        except Exception as e:
            logs.append(f"Mongo Handler Error: {e}")
            st.error(f"❌ Document handler failed: {e}")
            st.code(traceback.format_exc())

    # -------------------- POLICY HANDLER --------------------
    elif getattr(route, "value", str(route)).lower() == "policy":
        st.info("📜 Running Policy (RAG) handler...")
        logs.append("Router → Policy Handler (Mutlimedia)")
        try:
            if not get_cached_vectorstore():
                st.warning("⚠️ No in-memory Chroma DB found — please upload and embed documents first.")
                st.stop()

            with st.spinner("Executing Policy RAG retrieval..."):
                result_text = policy_handler(q)

            if isinstance(result_text, str) and result_text.startswith("ERROR"):
                st.error(result_text)
            else:
                st.success("✅ Policy handler executed successfully.")
        except Exception as e:
            logs.append(f"Policy Handler Error: {e}")
            st.error(f"❌ Policy handler failed: {e}")
            st.code(traceback.format_exc())

    # -------------------- BOTH HANDLERS --------------------
    elif getattr(route, "value", str(route)).lower() == "both":
        st.info("🔄 Running both Mongo & Policy handlers...")
        logs.append("Router → Both (Document + Policy)")
        try:
            from src.Mongo import query_mongo
            if not get_cached_vectorstore():
                st.warning("⚠️ No in-memory Chroma DB found — please upload and embed documents first.")
                st.stop()

            with st.spinner("Executing Document handler..."):
                res_doc = query_mongo(doc_q or q)
            with st.spinner("Executing Policy handler..."):
                res_pol = policy_handler(pol_q or q)

            if any(isinstance(x, str) and x.startswith("ERROR") for x in [res_doc, res_pol]):
                st.error("One or both handlers reported an error.")
            else:
                st.success("✅ Both handlers executed successfully.")

            result_text = f"--- DOCUMENT RESULT ---\n{res_doc}\n\n--- POLICY RESULT ---\n{res_pol}"
        except Exception as e:
            logs.append(f"Both Handler Error: {e}")
            st.error(f"❌ Combined handler failure: {e}")
            st.code(traceback.format_exc())

    else:
        st.warning(f"⚠️ Unknown route type: {route}")
        logs.append("Unknown Route")

    # -------------------- FINAL OUTPUT --------------------
    st.divider()
    st.subheader("🧾 Final Result")
    st.code(result_text or "No output generated.", language="text")

    st.subheader("📋 Execution Log")
    for log in logs:
        st.write("-", log)

    st.caption("✅ Process completed successfully — all stages executed inline.")
